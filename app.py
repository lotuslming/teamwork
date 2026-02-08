import os
import uuid
import json
import shutil
import time
import jwt as pyjwt
import requests
from datetime import datetime
from functools import wraps

from flask import Flask, request, jsonify, send_from_directory, render_template
from flask_jwt_extended import JWTManager, create_access_token, jwt_required, get_jwt_identity
from flask_socketio import SocketIO, emit, join_room, leave_room
from werkzeug.utils import secure_filename
import markdown
import bleach

from config import config
from models import db, User, Project, Card, Category, Attachment, ChatMessage, FileVersion, project_members, UnreadStatus

app = Flask(__name__)
env = os.environ.get('FLASK_ENV', 'development')
app.config.from_object(config.get(env, config['default']))

# Ensure upload folder exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(os.path.join(app.config['UPLOAD_FOLDER'], 'attachments'), exist_ok=True)
os.makedirs(os.path.join(app.config['UPLOAD_FOLDER'], 'chat'), exist_ok=True)

# Initialize extensions
db.init_app(app)
jwt = JWTManager(app)
socketio = SocketIO(app, cors_allowed_origins="*", async_mode='eventlet')

# Create tables
with app.app_context():
    db.create_all()

# Allowed file extensions
ALLOWED_EXTENSIONS = {'txt', 'md', 'pdf', 'doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx', 'png', 'jpg', 'jpeg', 'gif', 'zip', 'rar'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_file_type(filename):
    if '.' in filename:
        ext = filename.rsplit('.', 1)[1].lower()
        if ext in ['txt', 'md']:
            return 'text'
        elif ext in ['doc', 'docx']:
            return 'word'
        elif ext in ['xls', 'xlsx']:
            return 'excel'
        elif ext in ['ppt', 'pptx']:
            return 'powerpoint'
        elif ext in ['pdf']:
            return 'pdf'
        elif ext in ['png', 'jpg', 'jpeg', 'gif']:
            return 'image'
        else:
            return 'other'
    return 'other'

def extract_file_content(file_path, file_type, max_content_length=50000):
    """Extract text content from a file for full-text search.
    
    Returns extracted text, limited to max_content_length characters.
    Designed to fail gracefully and not block file operations.
    """
    try:
        content = None
        
        if file_type == 'text':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(max_content_length)
        
        elif file_type == 'word':
            from docx import Document
            doc = Document(file_path)
            content = '\n'.join([para.text for para in doc.paragraphs])
        
        elif file_type == 'excel':
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames[:5]:  # Limit to first 5 sheets
                ws = wb[sheet_name]
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    row_count += 1
                    if row_count > 1000:  # Limit rows per sheet
                        break
                    for cell in row:
                        if cell is not None:
                            text_parts.append(str(cell))
            wb.close()
            content = ' '.join(text_parts)
        
        elif file_type == 'powerpoint':
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides[:50]:  # Limit to first 50 slides
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text_parts.append(shape.text)
            content = '\n'.join(text_parts)
        
        elif file_type == 'pdf':
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            text_parts = []
            for page in reader.pages[:50]:  # Limit to first 50 pages
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            content = '\n'.join(text_parts)
        
        # Truncate to max length if needed
        if content and len(content) > max_content_length:
            content = content[:max_content_length]
        
        return content
    except Exception as e:
        # Log but don't fail - content indexing is optional
        print(f"Content extraction skipped for {file_path}: {e}")
        return None

# ================== AUTH ROUTES ==================

@app.route('/api/auth/register', methods=['POST'])
def register():
    data = request.get_json()
    
    if not data or not data.get('username') or not data.get('email') or not data.get('password'):
        return jsonify({'error': '请填写所有必填字段'}), 400
    
    if User.query.filter_by(username=data['username']).first():
        return jsonify({'error': '用户名已存在'}), 400
    
    if User.query.filter_by(email=data['email']).first():
        return jsonify({'error': '邮箱已被注册'}), 400
    
    user = User(
        username=data['username'],
        email=data['email'],
        avatar_color=data.get('avatar_color', '#00d4ff')
    )
    user.set_password(data['password'])
    
    db.session.add(user)
    db.session.commit()
    
    access_token = create_access_token(identity=str(user.id))
    return jsonify({
        'message': '注册成功',
        'user': user.to_dict(),
        'access_token': access_token
    }), 201

@app.route('/api/auth/login', methods=['POST'])
def login():
    data = request.get_json()
    
    if not data or not data.get('username') or not data.get('password'):
        return jsonify({'error': '请输入用户名和密码'}), 400
    
    user = User.query.filter_by(username=data['username']).first()
    
    if not user or not user.check_password(data['password']):
        return jsonify({'error': '用户名或密码错误'}), 401
    
    access_token = create_access_token(identity=str(user.id))
    return jsonify({
        'message': '登录成功',
        'user': user.to_dict(),
        'access_token': access_token
    })

@app.route('/api/auth/me', methods=['GET'])
@jwt_required()
def get_current_user():
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': '用户不存在'}), 404
    return jsonify(user.to_dict())

# ================== PROJECT ROUTES ==================

@app.route('/api/projects', methods=['GET'])
@jwt_required()
def get_projects():
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    
    # Get projects where user is owner or member
    projects = Project.query.filter(
        (Project.owner_id == user_id) | (Project.members.contains(user))
    ).order_by(Project.updated_at.desc()).all()
    
    result = []
    for p in projects:
        p_data = p.to_dict()
        # Calculate unread messages
        unread_stat = UnreadStatus.query.filter_by(user_id=user_id, project_id=p.id).first()
        last_read = unread_stat.last_read_at if unread_stat else datetime.min
        
        count = ChatMessage.query.filter(
            ChatMessage.project_id == p.id,
            ChatMessage.created_at > last_read,
            ChatMessage.user_id != user_id
        ).count()
        p_data['unread_count'] = count
        result.append(p_data)
        
    return jsonify(result)

@app.route('/api/projects/<int:project_id>/read', methods=['POST'])
@jwt_required()
def mark_project_read(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    status = UnreadStatus.query.filter_by(user_id=user_id, project_id=project_id).first()
    if not status:
        status = UnreadStatus(user_id=user_id, project_id=project_id)
        db.session.add(status)
    
    status.last_read_at = datetime.utcnow()
    db.session.commit()
    return jsonify({'success': True})

@app.route('/api/projects/<int:project_id>/columns', methods=['PUT'])
@jwt_required()
def update_layout(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if project.owner_id != user_id:
        return jsonify({'error': '只有项目所有者可以修改看板列'}), 403
        
    data = request.get_json()
    columns = data.get('columns')
    
    if not columns or not isinstance(columns, list):
         return jsonify({'error': '无效的列配置'}), 400
         
    # Optional: Check if deleted columns have cards (skipped per user request flexibility, or implemented if strict)
    # For now, allow modification. Cards in deleted columns will still have "column" string but wont show in board?
    # Or should we migrate them to TODO?
    # Better to prevent if cards exist.
    
    current_cols = set(project.columns or [])
    new_cols = set(columns)
    removed = current_cols - new_cols
    
    if removed:
        # Check if any cards in removed columns
        count = Card.query.filter(Card.project_id == project_id, Card.column.in_(removed)).count()
        if count > 0:
             return jsonify({'error': f'无法删除包含卡片的列: {", ".join(removed)}'}), 400

    project.columns = columns
    db.session.commit()
    return jsonify(project.to_dict())

@app.route('/api/projects', methods=['POST'])
@jwt_required()
def create_project():
    user_id = int(get_jwt_identity())
    data = request.get_json()
    
    if not data or not data.get('name'):
        return jsonify({'error': '请输入项目名称'}), 400
    
    project = Project(
        name=data['name'],
        description=data.get('description', ''),
        owner_id=user_id,
        columns=data.get('columns', ['待办', '进行中', '已完成'])
    )
    
    # Add owner as member
    owner = User.query.get(user_id)
    project.members.append(owner)
    
    db.session.add(project)
    db.session.commit()
    
    return jsonify(project.to_dict()), 201

@app.route('/api/projects/<int:project_id>', methods=['GET'])
@jwt_required()
def get_project(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    # Check if user has access
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权访问此项目'}), 403
    
    return jsonify(project.to_dict(include_cards=True))

@app.route('/api/projects/<int:project_id>', methods=['PUT'])
@jwt_required()
def update_project(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if project.owner_id != user_id:
        return jsonify({'error': '只有项目所有者可以编辑项目'}), 403
    
    data = request.get_json()
    
    if data.get('name'):
        project.name = data['name']
    if 'description' in data:
        project.description = data['description']
    if 'columns' in data:
        project.columns = data['columns']
    
    db.session.commit()
    return jsonify(project.to_dict())

@app.route('/api/projects/<int:project_id>', methods=['DELETE'])
@jwt_required()
def delete_project(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if project.owner_id != user_id:
        return jsonify({'error': '只有项目所有者可以删除项目'}), 403
    
    db.session.delete(project)
    db.session.commit()
    return jsonify({'message': '项目已删除'})

# ================== MEMBER ROUTES ==================

@app.route('/api/projects/<int:project_id>/invite', methods=['POST'])
@jwt_required()
def invite_member(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权邀请成员'}), 403
    
    data = request.get_json()
    username = data.get('username')
    
    if not username:
        return jsonify({'error': '请输入用户名'}), 400
    
    invitee = User.query.filter_by(username=username).first()
    if not invitee:
        return jsonify({'error': '用户不存在'}), 404
    
    if any(m.id == invitee.id for m in project.members):
        return jsonify({'error': '用户已经是项目成员'}), 400
    
    project.members.append(invitee)
    db.session.commit()
    
    return jsonify({'message': f'{username} 已加入项目', 'user': invitee.to_dict()})

@app.route('/api/projects/<int:project_id>/members', methods=['GET'])
@jwt_required()
def get_members(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权访问'}), 403
    
    return jsonify([m.to_dict() for m in project.members])

# ================== CARD ROUTES ==================

@app.route('/api/projects/<int:project_id>/cards', methods=['GET'])
@jwt_required()
def get_cards(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权访问'}), 403
    
    cards = Card.query.filter_by(project_id=project_id).order_by(Card.position).all()
    return jsonify([c.to_dict() for c in cards])

@app.route('/api/projects/<int:project_id>/cards', methods=['POST'])
@jwt_required()
def create_card(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权创建卡片'}), 403
    
    data = request.get_json()
    
    if not data or not data.get('title'):
        return jsonify({'error': '请输入卡片标题'}), 400
    
    # Get min position in column to add at top (min - 1)
    min_pos = db.session.query(db.func.min(Card.position)).filter_by(
        project_id=project_id,
        column=data.get('column', '待办')
    ).scalar()
    
    new_pos = (min_pos - 1) if min_pos is not None else 0
    
    card = Card(
        project_id=project_id,
        title=data['title'],
        content=data.get('content', ''),
        content_type=data.get('content_type', 'markdown'),
        column=data.get('column', '待办'),
        position=new_pos,
        due_date=datetime.fromisoformat(data['due_date']) if data.get('due_date') else None
    )
    
    # Add assignees
    if data.get('assignee_ids'):
        for uid in data['assignee_ids']:
            user = User.query.get(uid)
            if user and any(m.id == uid for m in project.members):
                card.assignees.append(user)
    
    # Add categories
    if data.get('category_ids'):
        for cid in data['category_ids']:
            category = Category.query.get(cid)
            if category and category.project_id == project_id:
                card.categories.append(category)
    
    db.session.add(card)
    db.session.commit()
    
    return jsonify(card.to_dict()), 201

@app.route('/api/cards/<int:card_id>', methods=['GET'])
@jwt_required()
def get_card(card_id):
    user_id = int(get_jwt_identity())
    card = Card.query.get_or_404(card_id)
    project = card.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权查看卡片'}), 403
    
    return jsonify(card.to_dict())

@app.route('/api/cards/<int:card_id>', methods=['PUT'])
@jwt_required()
def update_card(card_id):
    user_id = int(get_jwt_identity())
    card = Card.query.get_or_404(card_id)
    project = card.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权编辑卡片'}), 403
    
    data = request.get_json()
    
    if 'title' in data:
        card.title = data['title']
    if 'content' in data:
        card.content = data['content']
    if 'content_type' in data:
        card.content_type = data['content_type']
    if 'column' in data:
        card.column = data['column']
    if 'position' in data:
        card.position = data['position']
    if 'completed' in data:
        card.completed = data['completed']
    if 'due_date' in data:
        card.due_date = datetime.fromisoformat(data['due_date']) if data['due_date'] else None
    
    # Update assignees
    if 'assignee_ids' in data:
        card.assignees = []
        for uid in data['assignee_ids']:
            user = User.query.get(uid)
            if user and any(m.id == uid for m in project.members):
                card.assignees.append(user)
    
    # Update categories
    if 'category_ids' in data:
        card.categories = []
        for cid in data['category_ids']:
            category = Category.query.get(cid)
            if category and category.project_id == project.id:
                card.categories.append(category)
    
    db.session.commit()
    return jsonify(card.to_dict())

@app.route('/api/cards/<int:card_id>', methods=['DELETE'])
@jwt_required()
def delete_card(card_id):
    user_id = int(get_jwt_identity())
    card = Card.query.get_or_404(card_id)
    project = card.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权删除卡片'}), 403
    
    # Delete associated attachments files
    for attachment in card.attachments:
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'attachments', attachment.filename)
        if os.path.exists(file_path):
            os.remove(file_path)
    
    db.session.delete(card)
    db.session.commit()
    return jsonify({'message': '卡片已删除'})

@app.route('/api/cards/reorder', methods=['POST'])
@jwt_required()
def reorder_cards():
    user_id = int(get_jwt_identity())
    data = request.get_json()
    
    if not data or not data.get('cards'):
        return jsonify({'error': '无效数据'}), 400
    
    for card_data in data['cards']:
        card = Card.query.get(card_data['id'])
        if card:
            project = card.project
            if any(m.id == user_id for m in project.members):
                card.column = card_data.get('column', card.column)
                card.position = card_data.get('position', card.position)
    
    db.session.commit()
    return jsonify({'message': '卡片顺序已更新'})

# ================== SEARCH ROUTE ==================

@app.route('/api/projects/<int:project_id>/cards/search', methods=['GET'])
@jwt_required()
def search_cards(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权访问'}), 403
    
    query = Card.query.filter_by(project_id=project_id)
    
    # Text search
    q = request.args.get('q')
    include_attachments = request.args.get('include_attachments') == 'true'
    
    if q:
        if include_attachments:
            # Search in cards AND attachment filenames AND attachment contents
            query = query.outerjoin(Attachment).filter(
                (Card.title.ilike(f'%{q}%')) | 
                (Card.content.ilike(f'%{q}%')) |
                (Attachment.original_filename.ilike(f'%{q}%')) |
                (Attachment.content.ilike(f'%{q}%'))
            ).distinct()
        else:
            # Search only in cards
            query = query.filter(
                (Card.title.ilike(f'%{q}%')) | (Card.content.ilike(f'%{q}%'))
            )
    
    # Status filter
    status = request.args.get('status')
    if status == 'completed':
        query = query.filter(Card.completed == True)
    elif status == 'pending':
        query = query.filter(Card.completed == False)
    
    # Column filter
    column = request.args.get('column')
    if column:
        query = query.filter(Card.column == column)
    
    # Category filter
    category_id = request.args.get('category')
    if category_id:
        query = query.filter(Card.categories.any(id=int(category_id)))
    
    # Assignee filter
    assignee_id = request.args.get('assignee')
    if assignee_id:
        query = query.filter(Card.assignees.any(id=int(assignee_id)))
    
    cards = query.order_by(Card.position).all()
    return jsonify([c.to_dict() for c in cards])

# ================== CATEGORY ROUTES ==================

@app.route('/api/projects/<int:project_id>/categories', methods=['GET'])
@jwt_required()
def get_categories(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权访问'}), 403
    
    categories = Category.query.filter_by(project_id=project_id).all()
    return jsonify([c.to_dict() for c in categories])

@app.route('/api/projects/<int:project_id>/categories', methods=['POST'])
@jwt_required()
def create_category(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权创建类别'}), 403
    
    data = request.get_json()
    
    if not data or not data.get('name'):
        return jsonify({'error': '请输入类别名称'}), 400
    
    category = Category(
        project_id=project_id,
        name=data['name'],
        color=data.get('color', '#00d4ff')
    )
    
    db.session.add(category)
    db.session.commit()
    
    return jsonify(category.to_dict()), 201

@app.route('/api/categories/<int:category_id>', methods=['PUT'])
@jwt_required()
def update_category(category_id):
    user_id = int(get_jwt_identity())
    category = Category.query.get_or_404(category_id)
    project = category.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权编辑类别'}), 403
    
    data = request.get_json()
    
    if 'name' in data:
        category.name = data['name']
    if 'color' in data:
        category.color = data['color']
    
    db.session.commit()
    return jsonify(category.to_dict())

@app.route('/api/categories/<int:category_id>', methods=['DELETE'])
@jwt_required()
def delete_category(category_id):
    user_id = int(get_jwt_identity())
    category = Category.query.get_or_404(category_id)
    project = category.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权删除类别'}), 403
    
    db.session.delete(category)
    db.session.commit()
    return jsonify({'message': '类别已删除'})

# ================== ATTACHMENT ROUTES ==================

@app.route('/api/cards/<int:card_id>/attachments', methods=['POST'])
@jwt_required()
def upload_attachment(card_id):
    user_id = int(get_jwt_identity())
    card = Card.query.get_or_404(card_id)
    project = card.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权上传附件'}), 403
    
    if 'files' not in request.files:
        return jsonify({'error': '没有选择文件'}), 400
    
    files = request.files.getlist('files')
    attachments = []
    
    for file in files:
        if file and file.filename and allowed_file(file.filename):
            original_filename = secure_filename(file.filename)
            # Generate unique filename
            ext = original_filename.rsplit('.', 1)[1].lower() if '.' in original_filename else ''
            unique_filename = f"{uuid.uuid4().hex}.{ext}" if ext else uuid.uuid4().hex
            
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'attachments', unique_filename)
            file.save(file_path)
            
            # Extract text content from file for search
            file_type = get_file_type(original_filename)
            extracted_content = extract_file_content(file_path, file_type)
            
            attachment = Attachment(
                card_id=card_id,
                filename=unique_filename,
                original_filename=original_filename,
                file_type=file_type,
                file_size=os.path.getsize(file_path),
                content=extracted_content
            )
            db.session.add(attachment)
            attachments.append(attachment)
    
    db.session.commit()
    return jsonify([a.to_dict() for a in attachments]), 201

@app.route('/api/attachments/<int:attachment_id>', methods=['GET'])
@jwt_required()
def download_attachment(attachment_id):
    user_id = int(get_jwt_identity())
    attachment = Attachment.query.get_or_404(attachment_id)
    card = attachment.card
    project = card.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权下载附件'}), 403
    
    return send_from_directory(
        os.path.join(app.config['UPLOAD_FOLDER'], 'attachments'),
        attachment.filename,
        as_attachment=True,
        download_name=attachment.original_filename
    )

@app.route('/api/attachments/<int:attachment_id>/content', methods=['GET'])
@jwt_required()
def get_attachment_content(attachment_id):
    user_id = int(get_jwt_identity())
    attachment = Attachment.query.get_or_404(attachment_id)
    card = attachment.card
    project = card.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权访问附件'}), 403
    
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'attachments', attachment.filename)
    
    if attachment.file_type == 'text':
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        return jsonify({'content': content, 'type': 'text'})
    
    elif attachment.file_type == 'word':
        try:
            from docx import Document
            doc = Document(file_path)
            content = '\n\n'.join([para.text for para in doc.paragraphs])
            return jsonify({'content': content, 'type': 'word'})
        except Exception as e:
            return jsonify({'error': str(e)}), 500
    
    elif attachment.file_type == 'excel':
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            sheets = {}
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                data = []
                for row in ws.iter_rows(values_only=True):
                    data.append([str(cell) if cell is not None else '' for cell in row])
                sheets[sheet_name] = data
            return jsonify({'content': sheets, 'type': 'excel'})
        except Exception as e:
            return jsonify({'error': str(e)}), 500
    
    elif attachment.file_type == 'powerpoint':
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides = []
            for slide in prs.slides:
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_content.append(shape.text)
                slides.append('\n'.join(slide_content))
            return jsonify({'content': slides, 'type': 'powerpoint'})
        except Exception as e:
            return jsonify({'error': str(e)}), 500
    
    elif attachment.file_type == 'pdf':
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            pages = []
            for page in reader.pages:
                pages.append(page.extract_text() or '')
            return jsonify({'content': pages, 'type': 'pdf'})
        except Exception as e:
            return jsonify({'error': str(e)}), 500
    
    return jsonify({'error': '不支持的文件类型'}), 400

@app.route('/api/attachments/<int:attachment_id>/content', methods=['PUT'])
@jwt_required()
def update_attachment_content(attachment_id):
    user_id = int(get_jwt_identity())
    attachment = Attachment.query.get_or_404(attachment_id)
    card = attachment.card
    project = card.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权编辑附件'}), 403
    
    data = request.get_json()
    content = data.get('content')
    
    if content is None:
        return jsonify({'error': '没有内容'}), 400
    
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'attachments', attachment.filename)
    
    if attachment.file_type == 'text':
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(content)
        return jsonify({'message': '文件已保存'})
    
    elif attachment.file_type == 'word':
        try:
            from docx import Document
            doc = Document()
            for para in content.split('\n\n'):
                doc.add_paragraph(para)
            doc.save(file_path)
            return jsonify({'message': '文件已保存'})
        except Exception as e:
            return jsonify({'error': str(e)}), 500
    
    elif attachment.file_type == 'excel':
        try:
            import openpyxl
            wb = openpyxl.Workbook()
            for sheet_name, data in content.items():
                if sheet_name == list(content.keys())[0]:
                    ws = wb.active
                    ws.title = sheet_name
                else:
                    ws = wb.create_sheet(sheet_name)
                for row_idx, row in enumerate(data, 1):
                    for col_idx, value in enumerate(row, 1):
                        ws.cell(row=row_idx, column=col_idx, value=value)
            wb.save(file_path)
            return jsonify({'message': '文件已保存'})
        except Exception as e:
            return jsonify({'error': str(e)}), 500
    
    return jsonify({'error': '此文件类型不支持在线编辑'}), 400

@app.route('/api/attachments/<int:attachment_id>', methods=['DELETE'])
@jwt_required()
def delete_attachment(attachment_id):
    user_id = int(get_jwt_identity())
    attachment = Attachment.query.get_or_404(attachment_id)
    card = attachment.card
    project = card.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权删除附件'}), 403
    
    # Delete the physical file
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'attachments', attachment.filename)
    if os.path.exists(file_path):
        os.remove(file_path)
    
    # Also delete any version files
    versions_dir = os.path.join(app.config['UPLOAD_FOLDER'], 'versions')
    for version in attachment.versions.all():
        version_path = os.path.join(versions_dir, version.file_path)
        if os.path.exists(version_path):
            try:
                os.remove(version_path)
            except:
                pass  # Non-critical if version file delete fails
        db.session.delete(version)
    
    # Now delete the attachment record
    db.session.delete(attachment)
    db.session.commit()
    return jsonify({'message': '附件已删除'})

# ================== ONLYOFFICE INTEGRATION ==================

def generate_onlyoffice_token(payload):
    """Generate JWT token for OnlyOffice Document Server"""
    secret = app.config.get('ONLYOFFICE_JWT_SECRET', 'teamwork-onlyoffice-secret-key')
    return pyjwt.encode(payload, secret, algorithm='HS256')

def drop_onlyoffice_cache(doc_key):
    """Tell OnlyOffice to drop cached document so restored version is shown"""
    try:
        onlyoffice_url = app.config.get('ONLYOFFICE_URL', 'http://localhost:8080').rstrip('/')
        command_url = f"{onlyoffice_url}/coauthoring/CommandService.ashx"
        
        payload = {"c": "drop", "key": doc_key}
        
        # Add JWT token if configured
        secret = app.config.get('ONLYOFFICE_JWT_SECRET')
        if secret:
            payload["token"] = pyjwt.encode(payload, secret, algorithm='HS256')
        
        response = requests.post(command_url, json=payload, timeout=5)
        return response.status_code == 200
    except Exception as e:
        print(f"Failed to drop OnlyOffice cache for key {doc_key}: {e}")
        return False


@app.route('/api/attachments/<int:attachment_id>/onlyoffice-config', methods=['GET'])
@jwt_required()
def get_onlyoffice_config(attachment_id):
    """Generate OnlyOffice editor configuration"""
    user_id = int(get_jwt_identity())
    user = User.query.get_or_404(user_id)
    attachment = Attachment.query.get_or_404(attachment_id)
    card = attachment.card
    project = card.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权访问此附件'}), 403
    
    # Determine document type
    ext = attachment.original_filename.rsplit('.', 1)[-1].lower() if '.' in attachment.original_filename else ''
    doc_type_map = {
        'docx': 'word', 'doc': 'word',
        'xlsx': 'cell', 'xls': 'cell',
        'pptx': 'slide', 'ppt': 'slide'
    }
    document_type = doc_type_map.get(ext)
    
    if not document_type:
        return jsonify({'error': '此文件类型不支持OnlyOffice编辑'}), 400
    
    # Build file URL that OnlyOffice can access from Docker container
    # Use INTERNAL_URL config for Docker bridge access (not localhost which Docker can't reach)
    internal_url = app.config.get('INTERNAL_URL', 'http://172.17.0.1:5000').rstrip('/')
    file_url = f"{internal_url}/api/attachments/{attachment_id}/download"
    callback_url = f"{internal_url}/api/onlyoffice/callback"
    
    # Create a unique document key (changes when file is modified)
    # Use millisecond precision to ensure key changes even if modifications happen in same second
    doc_key = f"{attachment_id}_{int(attachment.uploaded_at.timestamp() * 1000)}"
    
    # Build OnlyOffice configuration
    config = {
        "document": {
            "fileType": ext,
            "key": doc_key,
            "title": attachment.original_filename,
            "url": file_url,
            "permissions": {
                "edit": True,
                "download": True,
                "print": True,
                "review": True,
                "comment": True
            }
        },
        "documentType": document_type,
        "editorConfig": {
            "mode": "edit",
            "lang": "zh-CN",
            "callbackUrl": callback_url,
            "user": {
                "id": str(user.id),
                "name": user.username
            },
            "customization": {
                "autosave": True,
                "forcesave": True,
                "chat": True,
                "comments": True,
                "compactHeader": False,
                "feedback": False,
                "help": False
            }
        },
        "token": ""
    }
    
    # Generate JWT token and add to config
    config["token"] = generate_onlyoffice_token(config)
    
    return jsonify({
        'config': config,
        'onlyoffice_url': app.config.get('ONLYOFFICE_URL', 'http://localhost:8080')
    })

@app.route('/api/attachments/<int:attachment_id>/download', methods=['GET'])
def download_attachment_for_onlyoffice(attachment_id):
    """Serve file for OnlyOffice to download"""
    attachment = Attachment.query.get_or_404(attachment_id)
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'attachments', attachment.filename)
    
    if not os.path.exists(file_path):
        return jsonify({'error': '文件不存在'}), 404
    
    return send_from_directory(
        os.path.join(app.config['UPLOAD_FOLDER'], 'attachments'),
        attachment.filename,
        as_attachment=True,
        download_name=attachment.original_filename
    )

@app.route('/api/onlyoffice/callback', methods=['POST'])
def onlyoffice_callback():
    """Handle OnlyOffice document save callback"""
    data = request.get_json()
    
    if not data:
        return jsonify({'error': 1}), 200
    
    status = data.get('status')
    
    # Status codes from OnlyOffice:
    # 0 - no document with the key identifier could be found
    # 1 - document is being edited
    # 2 - document is ready for saving (user closed editor)
    # 3 - document saving error has occurred
    # 4 - document is closed with no changes
    # 6 - document is being edited, but the current document state is saved
    # 7 - error has occurred while force saving the document
    
    if status in [2, 6]:  # Ready to save or force save
        try:
            download_url = data.get('url')
            key = data.get('key')
            users = data.get('users', [])
            
            if not download_url or not key:
                return jsonify({'error': 1}), 200
            
            # Parse attachment ID from key
            attachment_id = int(key.split('_')[0])
            attachment = Attachment.query.get(attachment_id)
            
            if not attachment:
                return jsonify({'error': 1}), 200
            
            # Create version backup before saving
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'attachments', attachment.filename)
            versions_dir = os.path.join(app.config['UPLOAD_FOLDER'], 'versions')
            os.makedirs(versions_dir, exist_ok=True)
            
            # Get current version number
            last_version = FileVersion.query.filter_by(attachment_id=attachment_id).order_by(FileVersion.version_number.desc()).first()
            new_version_num = (last_version.version_number + 1) if last_version else 1
            
            # Backup current file as version
            if os.path.exists(file_path):
                version_filename = f"{attachment_id}_v{new_version_num}_{attachment.filename}"
                version_path = os.path.join(versions_dir, version_filename)
                shutil.copy2(file_path, version_path)
                
                # Get user who made the edit
                editor_id = int(users[0]) if users else None
                
                # Create version record
                file_version = FileVersion(
                    attachment_id=attachment_id,
                    version_number=new_version_num,
                    file_path=version_filename,
                    file_size=os.path.getsize(version_path),
                    edited_by_id=editor_id or 1,
                    change_summary=f"Version {new_version_num} saved via OnlyOffice"
                )
                db.session.add(file_version)
            
            # Download the new file from OnlyOffice
            response = requests.get(download_url, timeout=30)
            if response.status_code == 200:
                with open(file_path, 'wb') as f:
                    f.write(response.content)
                
                # Update attachment metadata and commit immediately
                # This ensures the main save operation completes before content extraction
                attachment.file_size = len(response.content)
                attachment.uploaded_at = datetime.utcnow()
                
                try:
                    db.session.commit()
                except Exception as commit_error:
                    db.session.rollback()
                    print(f"OnlyOffice commit error: {commit_error}")
                    return jsonify({'error': 1}), 200
                
                # Re-index content for search AFTER commit (non-blocking, separate transaction)
                try:
                    extracted_content = extract_file_content(file_path, attachment.file_type)
                    if extracted_content:
                        attachment.content = extracted_content
                        db.session.commit()
                except Exception as extract_error:
                    db.session.rollback()
                    print(f"Content extraction error (non-critical): {extract_error}")
            
            return jsonify({'error': 0}), 200
        except Exception as e:
            db.session.rollback()
            print(f"OnlyOffice callback error: {e}")
            return jsonify({'error': 1}), 200
    
    return jsonify({'error': 0}), 200

@app.route('/api/attachments/<int:attachment_id>/versions', methods=['GET'])
@jwt_required()
def get_file_versions(attachment_id):
    """Get version history for an attachment"""
    user_id = int(get_jwt_identity())
    attachment = Attachment.query.get_or_404(attachment_id)
    card = attachment.card
    project = card.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权访问版本历史'}), 403
    
    versions = FileVersion.query.filter_by(attachment_id=attachment_id).order_by(FileVersion.version_number.desc()).all()
    return jsonify({'versions': [v.to_dict() for v in versions]})

@app.route('/api/attachments/<int:attachment_id>/restore/<int:version_id>', methods=['POST'])
@jwt_required()
def restore_file_version(attachment_id, version_id):
    """Restore a previous version of an attachment"""
    user_id = int(get_jwt_identity())
    attachment = Attachment.query.get_or_404(attachment_id)
    card = attachment.card
    project = card.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权恢复版本'}), 403
    
    version = FileVersion.query.get_or_404(version_id)
    if version.attachment_id != attachment_id:
        return jsonify({'error': '版本不匹配'}), 400
    
    versions_dir = os.path.join(app.config['UPLOAD_FOLDER'], 'versions')
    version_path = os.path.join(versions_dir, version.file_path)
    current_path = os.path.join(app.config['UPLOAD_FOLDER'], 'attachments', attachment.filename)
    
    if not os.path.exists(version_path):
        return jsonify({'error': '版本文件不存在'}), 404
    
    # Calculate the OLD document key before restore (to invalidate OnlyOffice cache)
    old_doc_key = f"{attachment_id}_{int(attachment.uploaded_at.timestamp() * 1000)}"
    
    # Create a new version of current file before restoring
    last_version = FileVersion.query.filter_by(attachment_id=attachment_id).order_by(FileVersion.version_number.desc()).first()
    new_version_num = (last_version.version_number + 1) if last_version else 1
    
    if os.path.exists(current_path):
        backup_filename = f"{attachment_id}_v{new_version_num}_{attachment.filename}"
        backup_path = os.path.join(versions_dir, backup_filename)
        shutil.copy2(current_path, backup_path)
        
        backup_version = FileVersion(
            attachment_id=attachment_id,
            version_number=new_version_num,
            file_path=backup_filename,
            file_size=os.path.getsize(backup_path),
            edited_by_id=user_id,
            change_summary=f"Backup before restoring to version {version.version_number}"
        )
        db.session.add(backup_version)
    
    # Restore the old version
    shutil.copy2(version_path, current_path)
    attachment.file_size = os.path.getsize(current_path)
    attachment.uploaded_at = datetime.utcnow()
    db.session.commit()
    
    # Drop OnlyOffice cache so the restored version is shown
    drop_onlyoffice_cache(old_doc_key)
    
    # Calculate new doc key for frontend
    new_doc_key = f"{attachment_id}_{int(attachment.uploaded_at.timestamp() * 1000)}"
    
    return jsonify({
        'message': f'已恢复到版本 {version.version_number}',
        'attachment': attachment.to_dict(),
        'new_doc_key': new_doc_key,
        'cache_dropped': True
    })

@app.route('/api/attachments/<int:attachment_id>/save-version', methods=['POST'])
@jwt_required()
def save_manual_version(attachment_id):
    """Manually save a version of the current file (independent of OnlyOffice)"""
    user_id = int(get_jwt_identity())
    attachment = Attachment.query.get_or_404(attachment_id)
    card = attachment.card
    project = card.project
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权保存版本'}), 403
    
    data = request.get_json() or {}
    change_summary = data.get('summary', '手动保存版本')
    
    current_path = os.path.join(app.config['UPLOAD_FOLDER'], 'attachments', attachment.filename)
    
    if not os.path.exists(current_path):
        return jsonify({'error': '文件不存在'}), 404
    
    versions_dir = os.path.join(app.config['UPLOAD_FOLDER'], 'versions')
    os.makedirs(versions_dir, exist_ok=True)
    
    # Get next version number
    last_version = FileVersion.query.filter_by(attachment_id=attachment_id).order_by(FileVersion.version_number.desc()).first()
    new_version_num = (last_version.version_number + 1) if last_version else 1
    
    # Create version backup
    version_filename = f"{attachment_id}_v{new_version_num}_{attachment.filename}"
    version_path = os.path.join(versions_dir, version_filename)
    shutil.copy2(current_path, version_path)
    
    # Create version record
    file_version = FileVersion(
        attachment_id=attachment_id,
        version_number=new_version_num,
        file_path=version_filename,
        file_size=os.path.getsize(version_path),
        edited_by_id=user_id,
        change_summary=change_summary
    )
    db.session.add(file_version)
    db.session.commit()
    
    return jsonify({
        'message': f'已保存版本 {new_version_num}',
        'version': file_version.to_dict()
    })

@app.route('/api/projects/<int:project_id>/ai/ask', methods=['POST'])
@jwt_required()
def ai_ask(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权访问'}), 403
    
    data = request.get_json()
    question = data.get('question')
    
    if not question:
        return jsonify({'error': '请输入问题'}), 400
    
    # Build context from project cards
    cards = Card.query.filter_by(project_id=project_id).all()
    context = f"项目名称: {project.name}\n项目描述: {project.description}\n\n卡片内容:\n"
    for card in cards:
        context += f"- [{card.column}] {card.title}: {card.content[:500]}\n"
    
    api_key = app.config.get('OPENAI_API_KEY')
    if not api_key:
        return jsonify({'error': '请配置OpenAI API密钥'}), 400
    
    try:
        from openai import OpenAI
        client = OpenAI(api_key=api_key)
        
        response = client.chat.completions.create(
            model=app.config.get('OPENAI_MODEL', 'gpt-3.5-turbo'),
            messages=[
                {"role": "system", "content": f"你是一个项目助手。根据以下项目信息回答用户的问题:\n\n{context}"},
                {"role": "user", "content": question}
            ],
            max_tokens=1000
        )
        
        answer = response.choices[0].message.content
        return jsonify({'answer': answer})
    
    except Exception as e:
        return jsonify({'error': f'AI请求失败: {str(e)}'}), 500

@app.route('/api/projects/<int:project_id>/ai/summarize', methods=['POST'])
@jwt_required()
def ai_summarize(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权访问'}), 403
    
    data = request.get_json()
    card_ids = data.get('card_ids', [])
    
    if not card_ids:
        return jsonify({'error': '请选择要总结的卡片'}), 400
    
    cards = Card.query.filter(Card.id.in_(card_ids), Card.project_id == project_id).all()
    
    if not cards:
        return jsonify({'error': '未找到指定卡片'}), 404
    
    content = "\n\n".join([f"## {card.title}\n{card.content}" for card in cards])
    
    api_key = app.config.get('OPENAI_API_KEY')
    if not api_key:
        return jsonify({'error': '请配置OpenAI API密钥'}), 400
    
    try:
        from openai import OpenAI
        client = OpenAI(api_key=api_key)
        
        response = client.chat.completions.create(
            model=app.config.get('OPENAI_MODEL', 'gpt-3.5-turbo'),
            messages=[
                {"role": "system", "content": "请将以下内容总结成一份文档，保持结构清晰、内容完整。"},
                {"role": "user", "content": content}
            ],
            max_tokens=2000
        )
        
        summary = response.choices[0].message.content
        return jsonify({'summary': summary})
    
    except Exception as e:
        return jsonify({'error': f'AI请求失败: {str(e)}'}), 500

@app.route('/api/ai/config', methods=['PUT'])
@jwt_required()
def update_ai_config():
    data = request.get_json()
    
    if data.get('api_base'):
        app.config['OPENAI_API_BASE'] = data['api_base']
    if data.get('api_key'):
        app.config['OPENAI_API_KEY'] = data['api_key']
    if data.get('model'):
        app.config['OPENAI_MODEL'] = data['model']
    
    return jsonify({'message': 'AI配置已更新'})

# ================== CHAT ROUTES ==================

@app.route('/api/projects/<int:project_id>/messages', methods=['GET'])
@jwt_required()
def get_messages(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权访问'}), 403
    
    # Pagination
    page = request.args.get('page', 1, type=int)
    per_page = request.args.get('per_page', 50, type=int)
    
    messages = ChatMessage.query.filter_by(project_id=project_id)\
        .order_by(ChatMessage.created_at.desc())\
        .paginate(page=page, per_page=per_page, error_out=False)
    
    return jsonify({
        'messages': [m.to_dict() for m in reversed(messages.items)],
        'has_more': messages.has_next,
        'total': messages.total
    })

@app.route('/api/projects/<int:project_id>/messages', methods=['POST'])
@jwt_required()
def post_message(project_id):
    user_id = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)
    
    if not any(m.id == user_id for m in project.members):
        return jsonify({'error': '无权发送消息'}), 403
    
    content = request.form.get('content', '')
    file = request.files.get('file')
    
    if not content and not file:
        return jsonify({'error': '请输入消息内容或上传文件'}), 400
    
    file_path = None
    file_name = None
    
    if file and file.filename:
        original_filename = secure_filename(file.filename)
        ext = original_filename.rsplit('.', 1)[1].lower() if '.' in original_filename else ''
        unique_filename = f"{uuid.uuid4().hex}.{ext}" if ext else uuid.uuid4().hex
        
        save_path = os.path.join(app.config['UPLOAD_FOLDER'], 'chat', unique_filename)
        file.save(save_path)
        file_path = unique_filename
        file_name = original_filename
    
    message = ChatMessage(
        project_id=project_id,
        user_id=user_id,
        content=content,
        file_path=file_path,
        file_name=file_name
    )
    
    db.session.add(message)
    db.session.commit()
    
    # Emit to socket room
    socketio.emit('new_message', message.to_dict(), room=f'project_{project_id}')
    
    return jsonify(message.to_dict()), 201

@app.route('/api/chat/files/<filename>', methods=['GET'])
def get_chat_file(filename):
    """Serve chat files - no auth required for direct file access"""
    return send_from_directory(
        os.path.join(app.config['UPLOAD_FOLDER'], 'chat'),
        filename
    )

@app.route('/api/chat/files/<filename>/onlyoffice-config', methods=['GET'])
@jwt_required()
def get_chat_file_onlyoffice_config(filename):
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'chat', filename)
    if not os.path.exists(file_path):
        return jsonify({'error': '文件不存在'}), 404
        
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    doc_type_map = {
        'docx': 'word', 'doc': 'word',
        'xlsx': 'cell', 'xls': 'cell',
        'pptx': 'slide', 'ppt': 'slide'
    }
    document_type = doc_type_map.get(ext)
    if not document_type:
         return jsonify({'error': '不支持此类文件'}), 400

    internal_url = app.config.get('INTERNAL_URL', 'http://172.17.0.1:5000').rstrip('/')
    file_url = f"{internal_url}/api/chat/files/{filename}"
    callback_url = f"{internal_url}/api/onlyoffice/chat/callback"
    
    mtime = int(os.path.getmtime(file_path) * 1000)
    doc_key = f"chat_{filename}_{mtime}"
    
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)

    config = {
        "document": {
            "fileType": ext,
            "key": doc_key,
            "title": filename,
            "url": file_url,
            "permissions": {
                "edit": True,
                "download": True,
            }
        },
        "documentType": document_type,
        "editorConfig": {
            "callbackUrl": callback_url,
            "user": {"id": str(user.id), "name": user.username},
            "lang": "zh-CN"
        }
    }
    
    msg = ChatMessage.query.filter_by(file_path=filename).first()
    if msg:
        config['document']['title'] = msg.file_name
        
    config['token'] = generate_onlyoffice_token(config)
    
    return jsonify({
        'config': config,
        'onlyoffice_url': app.config.get('ONLYOFFICE_URL', 'http://localhost:8080')
    })

@app.route('/api/onlyoffice/chat/callback', methods=['POST'])
def onlyoffice_chat_callback():
    data = request.get_json()
    if not data: return jsonify({'error': 1})
    
    status = data.get('status')
    if status in [2, 6]:
        download_url = data.get('url')
        key = data.get('key')
        
        if key and key.startswith('chat_'):
             parts = key.split('_')
             if len(parts) >= 2:
                 filename = parts[1]
                 file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'chat', filename)
                 
                 try:
                     response = requests.get(download_url, timeout=30)
                     if response.status_code == 200:
                         with open(file_path, 'wb') as f:
                             f.write(response.content)
                         return jsonify({'error': 0})
                 except Exception as e:
                     print(f"Chat callback error: {e}")
                     return jsonify({'error': 1})
                 
    return jsonify({'error': 0})

# ================== WEBSOCKET HANDLERS ==================

@socketio.on('join')
def on_join(data):
    room = f"project_{data['project_id']}"
    join_room(room)
    emit('user_joined', {'user_id': data.get('user_id')}, room=room)

@socketio.on('leave')
def on_leave(data):
    room = f"project_{data['project_id']}"
    leave_room(room)
    emit('user_left', {'user_id': data.get('user_id')}, room=room)

@socketio.on('typing')
def on_typing(data):
    room = f"project_{data['project_id']}"
    emit('user_typing', {'user_id': data.get('user_id'), 'username': data.get('username')}, room=room, include_self=False)

# ================== MARKDOWN HELPER ==================

@app.route('/api/render-markdown', methods=['POST'])
def render_markdown():
    data = request.get_json()
    content = data.get('content', '')
    
    # Convert markdown to HTML
    html = markdown.markdown(content, extensions=['tables', 'fenced_code', 'codehilite'])
    
    # Sanitize HTML
    allowed_tags = ['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'br', 'hr', 
                    'strong', 'em', 'a', 'ul', 'ol', 'li', 'code', 'pre',
                    'blockquote', 'table', 'thead', 'tbody', 'tr', 'th', 'td',
                    'img', 'span', 'div']
    allowed_attrs = {'a': ['href', 'title'], 'img': ['src', 'alt'], '*': ['class']}
    
    clean_html = bleach.clean(html, tags=allowed_tags, attributes=allowed_attrs)
    
    return jsonify({'html': clean_html})

# ================== HEALTH CHECK ENDPOINTS ==================

@app.route('/health')
def health_check():
    """Health check endpoint for load balancers and monitoring.
    Returns 200 if the application is running.
    """
    return jsonify({
        'status': 'healthy',
        'timestamp': datetime.utcnow().isoformat(),
        'service': 'teamwork'
    }), 200

@app.route('/ready')
def readiness_check():
    """Readiness check endpoint that verifies database connectivity.
    Returns 200 if ready to serve requests, 503 if not.
    """
    try:
        # Test database connection
        db.session.execute(db.text('SELECT 1'))
        db_status = 'connected'
    except Exception as e:
        return jsonify({
            'status': 'not_ready',
            'database': 'disconnected',
            'error': str(e)
        }), 503
    
    return jsonify({
        'status': 'ready',
        'database': db_status,
        'timestamp': datetime.utcnow().isoformat()
    }), 200

# ================== MAIN PAGE ==================

@app.route('/')
def index():
    return render_template('index.html')

if __name__ == '__main__':
    socketio.run(app, debug=True, host='0.0.0.0', port=5000)
